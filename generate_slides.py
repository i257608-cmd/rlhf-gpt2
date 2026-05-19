"""
Professional RLHF Presentation Generator
==========================================
Generates: results/day10/RLHF_Presentation.pptx
Authors: Muhammad Haseeb, Muavia Shakeel — FAST University Islamabad
"""

import os, json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

OUT_DIR = "results/day10"
os.makedirs(OUT_DIR, exist_ok=True)

# ── Load results ──────────────────────────────────────────────
c1   = json.load(open("results/c1/eval_results.json"))
c2   = json.load(open("results/c2/comparison.json"))
d8   = json.load(open("results/day8/alignment_tax.json"))
sweep = pd.read_csv("results/stability_sweep/sweep_results.csv")

agg    = c1["aggregate"]
sft_c2 = c2["SFT (baseline)"]
ppo_c2 = c2["PPO (aligned)"]
dpo_c2 = c2["DPO (aligned)"]
opt    = d8["optimal_stopping"]
method = d8["method_alignment_tax"]

# Sweep-derived C1 averages (3 seeds, β=0.20) — must match Slide 6 table
ppo_b02       = sweep[(sweep["beta"] == 0.2) & (sweep["status"] == "ok")]
c1_before_mean = ppo_b02["mean_reward_before"].mean()
c1_after_mean  = ppo_b02["mean_reward_after"].mean()
c1_delta_mean  = c1_after_mean - c1_before_mean
c1_improv_mean = ppo_b02["reward_improvement_pct"].mean()

# ── Design constants ──────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Colour palette
C_DARK_BLUE  = RGBColor(0x1A, 0x37, 0x6B)   # dark navy
C_MID_BLUE   = RGBColor(0x2E, 0x6D, 0xC8)   # accent blue
C_LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF7)   # background tint
C_ORANGE     = RGBColor(0xF0, 0x7D, 0x00)   # accent highlight
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x1E, 0x1E, 0x1E)
C_GREY       = RGBColor(0x60, 0x60, 0x60)
C_GREEN      = RGBColor(0x1A, 0x7A, 0x3C)
C_RED        = RGBColor(0xC0, 0x00, 0x00)

# Figure paths
FIG_REWARD = "results/c1/reward_curve.png"
FIG_RVK    = "results/stability_sweep/reward_vs_kl.png"
FIG_BETA   = "results/stability_sweep/beta_summary.png"
FIG_PPL    = "results/day8/ppl_vs_reward.png"
FIG_TAX    = "results/day8/tax_per_beta.png"
FIG_OPT    = "results/day8/optimal_kl.png"

# ── Presentation object ───────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # fully blank layout

# ══════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════
def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width     = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h, font_name="Calibri", font_size=18,
                 bold=False, italic=False, color=C_BLACK, align=PP_ALIGN.LEFT,
                 word_wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font_name
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_text_box_multi(slide, lines, l, t, w, h, font_name="Calibri",
                       font_size=18, bold=False, color=C_BLACK,
                       align=PP_ALIGN.LEFT, line_spacing_pt=None):
    """lines: list of (text, bold, size, color)"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            item = (item, bold, font_size, color)
        txt, is_bold, size, clr = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.name  = font_name
        run.font.size  = Pt(size)
        run.font.bold  = is_bold
        run.font.color.rgb = clr
        if line_spacing_pt:
            from pptx.util import Pt as oPt
            from pptx.oxml.ns import qn
            from lxml import etree
    return txBox

def add_slide_header(slide, title, subtitle=None):
    """Top navy bar + title text."""
    bar_h = Inches(1.1)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, fill_rgb=C_DARK_BLUE)
    # accent line
    add_rect(slide, 0, bar_h, SLIDE_W, Pt(4), fill_rgb=C_ORANGE)

    add_text_box(slide, title, Inches(0.4), Inches(0.1),
                 Inches(11.5), Inches(0.8),
                 font_name="Calibri", font_size=32, bold=True,
                 color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, Inches(0.4), Inches(0.75),
                     Inches(11), Inches(0.35),
                     font_name="Calibri", font_size=16, italic=True,
                     color=RGBColor(0xC0, 0xD8, 0xF8), align=PP_ALIGN.LEFT)

def add_slide_footer(slide, slide_num, total=13):
    foot_y = Inches(7.1)
    foot_h = Inches(0.4)
    add_rect(slide, 0, foot_y, SLIDE_W, foot_h, fill_rgb=C_DARK_BLUE)
    add_text_box(slide,
                 "FAST National University | Islamabad | RLHF on GPT-2",
                 Inches(0.3), foot_y, Inches(10), foot_h,
                 font_size=10, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, f"{slide_num}/{total}",
                 Inches(12.3), foot_y, Inches(0.8), foot_h,
                 font_size=10, color=C_WHITE, align=PP_ALIGN.RIGHT)

def add_bullet_box(slide, items, l, t, w, h, font_size=16, indent_size=14):
    """items: list of (text, level) where level 0=main, 1=sub"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            item = (item, 0)
        txt, level = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if level == 0:
            p.add_run().text = "▸  " + txt
        else:
            p.add_run().text = "     ◦  " + txt
        sz = font_size if level == 0 else indent_size
        p.runs[0].font.name  = "Calibri"
        p.runs[0].font.size  = Pt(sz)
        p.runs[0].font.color.rgb = C_BLACK if level == 0 else C_GREY
        p.alignment = PP_ALIGN.LEFT
    return txBox

def add_table(slide, headers, rows, l, t, w, h, header_fill=C_DARK_BLUE,
              font_size=13, alt_fill=C_LIGHT_BLUE):
    from pptx.util import Inches, Pt
    rows_total = 1 + len(rows)
    table = slide.shapes.add_table(rows_total, len(headers), l, t, w, h).table
    # Column widths equal
    col_w = int(w / len(headers))
    for ci in range(len(headers)):
        table.columns[ci].width = col_w
    # Header
    for ci, h_txt in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h_txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name  = "Calibri"
        run.font.size  = Pt(font_size)
        run.font.bold  = True
        run.font.color.rgb = C_WHITE
    # Data
    for ri, row_data in enumerate(rows):
        bg = alt_fill if ri % 2 == 1 else C_WHITE
        for ci, val in enumerate(row_data):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.name  = "Calibri"
            run.font.size  = Pt(font_size)
            run.font.color.rgb = C_BLACK
    return table

def add_image_safe(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        add_rect(slide, l, t, w, h, fill_rgb=C_LIGHT_BLUE,
                 line_rgb=C_MID_BLUE, line_width=Pt(1))
        add_text_box(slide, f"[Figure: {os.path.basename(path)}]",
                     l + Inches(0.1), t + h/2 - Inches(0.2), w - Inches(0.2), Inches(0.4),
                     font_size=12, color=C_GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
# Background gradient (dark blue to lighter)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_DARK_BLUE)
# Side accent bar
add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_rgb=C_ORANGE)
# Bottom accent
add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.4), fill_rgb=C_MID_BLUE)

# Main title
add_text_box(slide,
    "Reinforcement Learning from Human Feedback on GPT-2",
    Inches(0.5), Inches(1.5), Inches(12), Inches(1.2),
    font_size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide,
    "A Comparative Study of PPO, DPO, and KL-Penalty Stability",
    Inches(0.5), Inches(2.7), Inches(12), Inches(0.8),
    font_size=22, italic=True,
    color=RGBColor(0xA0, 0xC4, 0xFF), align=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Inches(2), Inches(3.65), Inches(9), Pt(3), fill_rgb=C_ORANGE)

# Authors
add_text_box(slide,
    "Muhammad Haseeb   |   Muavia Shakeel",
    Inches(0.5), Inches(3.9), Inches(12), Inches(0.6),
    font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide,
    "Department of Computer Science\nFAST National University of Computer and Emerging Sciences, Islamabad",
    Inches(0.5), Inches(4.55), Inches(12), Inches(0.9),
    font_size=16, color=RGBColor(0xC0, 0xD8, 0xF8), align=PP_ALIGN.CENTER)

add_text_box(slide, "Masters in AI  |  Semester Project",
    Inches(0.5), Inches(5.6), Inches(12), Inches(0.5),
    font_size=14, italic=True, color=RGBColor(0x90, 0xAA, 0xCC),
    align=PP_ALIGN.CENTER)

add_slide_footer(slide, 1)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "Presentation Outline")

agenda = [
    ("Background & Motivation", "Why RLHF? The alignment problem"),
    ("Three Research Contributions", "C1: PPO Pipeline  •  C2: PPO vs DPO  •  C3: KL Stability"),
    ("Methodology", "SFT → Reward Model → PPO / DPO → Evaluation"),
    ("Reward Model Results", "85.5% accuracy proxy reward"),
    ("C1 Results: PPO Training", "Per-seed reward improvement"),
    ("C2 Results: PPO vs DPO", "Reward efficiency comparison"),
    ("C3 Results: Stability Sweep", "Hacking onset & optimal β"),
    ("Alignment Tax Analysis", "Perplexity trade-off curves"),
    ("Key Findings & Conclusion", "Recommendations"),
]

for i, (title, sub) in enumerate(agenda):
    y = Inches(1.3) + i * Inches(0.58)
    # Number badge
    add_rect(slide, Inches(0.4), y, Inches(0.45), Inches(0.42),
             fill_rgb=C_MID_BLUE)
    add_text_box(slide, str(i+1),
                 Inches(0.4), y, Inches(0.45), Inches(0.42),
                 font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title,
                 Inches(1.05), y, Inches(5.5), Inches(0.44),
                 font_size=16, bold=True, color=C_DARK_BLUE)
    add_text_box(slide, sub,
                 Inches(6.8), y, Inches(6), Inches(0.44),
                 font_size=14, color=C_GREY)

add_slide_footer(slide, 2)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Background & Motivation
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "Background & Motivation",
                 "The alignment problem in language models")

bullets = [
    ("The Alignment Problem", 0),
    ("LLMs trained on next-token prediction produce fluent but misaligned text", 1),
    ("GPT-2 may generate harmful, biased, or off-topic content by default", 1),
    ("RLHF: Solution Approach", 0),
    ("Train a reward model from human preference labels", 1),
    ("Fine-tune the LM policy to maximise the reward signal (PPO)", 1),
    ("Key Challenge: Reward Hacking", 0),
    ("Over-optimising a proxy reward diverges from true human preferences", 1),
    ("KL divergence penalty required to prevent policy collapse", 1),
    ("Research Questions", 0),
    ("Can PPO reliably improve GPT-2 alignment at small scale?", 1),
    ("Does PPO outperform DPO in reward efficiency?", 1),
    ("What KL threshold separates stable training from reward hacking?", 1),
]
add_bullet_box(slide, bullets, Inches(0.5), Inches(1.35), Inches(12.5), Inches(5.5))
add_slide_footer(slide, 3)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Pipeline Overview
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "Methodology: RLHF Pipeline",
                 "GPT-2 (124M) + IMDB Dataset")

stages = [
    ("1\nSFT", "Supervised Fine-Tuning\nIMDB reviews\n128-token context"),
    ("2\nRM",  "Reward Model\n85.5% acc.\nGPT-2 + classifier head"),
    ("3a\nPPO","PPO Training\nβ∈{0.05–0.50}\n200 steps × 3 seeds"),
    ("3b\nDPO","DPO Training\nPreference pairs\nβ=0.10 × 3 seeds"),
    ("4\nEval","Evaluation\nReward · PPO · KL\nAlignment tax"),
]
colors = [C_MID_BLUE, C_MID_BLUE, C_DARK_BLUE, RGBColor(0x5A, 0x3A, 0xA0), C_GREEN]

box_w = Inches(2.2)
box_h = Inches(2.8)
gap   = Inches(0.25)
start_x = Inches(0.5)
y     = Inches(2.0)

for i, (label, desc) in enumerate(stages):
    x = start_x + i * (box_w + gap)
    add_rect(slide, x, y, box_w, box_h, fill_rgb=colors[i])
    # Label
    add_text_box(slide, label, x, y + Inches(0.1),
                 box_w, Inches(0.8),
                 font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # desc
    add_text_box(slide, desc, x, y + Inches(0.9),
                 box_w, Inches(1.8),
                 font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)
    # Arrow
    if i < len(stages) - 1:
        ax = x + box_w + Inches(0.02)
        add_text_box(slide, "→",
                     ax, y + Inches(1.1), gap + Inches(0.01), Inches(0.6),
                     font_size=24, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# Dataset box at bottom
add_rect(slide, Inches(3.5), Inches(5.15), Inches(6.3), Inches(0.8),
         fill_rgb=C_LIGHT_BLUE, line_rgb=C_MID_BLUE, line_width=Pt(1))
add_text_box(slide,
    "Dataset: IMDB (50k reviews)  |  Model: GPT-2 124M  |  Framework: HuggingFace TRL 0.24",
    Inches(3.5), Inches(5.15), Inches(6.3), Inches(0.8),
    font_size=13, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

add_slide_footer(slide, 4)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Reward Model Results
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "Reward Model Performance",
                 "GPT-2 + Classification Head — 3 Epochs on IMDB")

# Metric cards
metrics = [
    ("85.5%",  "Evaluation Accuracy", C_GREEN),
    ("0.356",  "Evaluation Loss",     C_MID_BLUE),
    ("3",      "Training Epochs",     C_DARK_BLUE),
    ("GPT-2",  "Base Architecture",   C_ORANGE),
]
for i, (val, label, color) in enumerate(metrics):
    x = Inches(0.5) + i * Inches(3.1)
    y_c = Inches(2.0)
    add_rect(slide, x, y_c, Inches(2.8), Inches(2.2), fill_rgb=color)
    add_text_box(slide, val, x, y_c + Inches(0.3),
                 Inches(2.8), Inches(1.0),
                 font_size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, x, y_c + Inches(1.35),
                 Inches(2.8), Inches(0.6),
                 font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide,
    "The reward model reliably distinguishes positive from negative IMDB reviews. "
    "Frozen at 85.5% accuracy, it serves as a stable proxy reward signal for all "
    "downstream PPO and DPO training runs.",
    Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.2),
    font_size=16, color=C_DARK_BLUE, align=PP_ALIGN.LEFT)

add_slide_footer(slide, 5)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — C1: PPO Per-Seed Results
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "C1: PPO Training Results",
                 "β = 0.20 · 200 Steps · 3 Seeds · No Reward Hacking")

# Table — pull per-seed data from stability sweep at beta=0.20
ppo_seeds = sweep[(sweep["beta"] == 0.2) & (sweep["status"] == "ok")]
rows = []
for r in ppo_seeds.itertuples():
    rows.append([
        f"Seed {int(r.seed)}",
        f"{r.mean_reward_before:.3f}",
        f"{r.mean_reward_after:.3f}",
        f"+{r.mean_reward_after - r.mean_reward_before:.3f}",
        f"+{r.reward_improvement_pct:.1f}%",
        f"{r.max_kl:.2f}",
        "✓ Stable"
    ])

add_table(slide,
    ["Seed", "Before (SFT)", "After (PPO)", "Delta", "Improvement", "Max KL", "Status"],
    rows,
    Inches(0.4), Inches(1.4), Inches(7.8), Inches(1.8))

# Image
add_image_safe(slide, FIG_REWARD,
               Inches(8.4), Inches(1.3), Inches(4.6), Inches(3.4))
add_text_box(slide, "Fig. 1 — Reward learning curve (seed=42)",
             Inches(8.4), Inches(4.75), Inches(4.6), Inches(0.4),
             font_size=11, italic=True, color=C_GREY, align=PP_ALIGN.CENTER)

# Summary box
add_rect(slide, Inches(0.4), Inches(3.4), Inches(7.8), Inches(1.3),
         fill_rgb=C_LIGHT_BLUE, line_rgb=C_MID_BLUE, line_width=Pt(1))
add_text_box(slide,
    f"Training Reward:  {c1_before_mean:.3f}  →  {c1_after_mean:.3f}  "
    f"(Δ = +{c1_delta_mean:.3f}, +{c1_improv_mean:.1f}%)     |     Post-train Eval Alignment Tax: +{agg['alignment_tax_pct']:.1f}% PPL",
    Inches(0.5), Inches(3.5), Inches(7.5), Inches(0.55),
    font_size=15, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide,
    "No reward collapse detected across any seed. All max KL values < 15.",
    Inches(0.5), Inches(4.0), Inches(7.5), Inches(0.5),
    font_size=14, color=C_GREY, align=PP_ALIGN.CENTER)

add_slide_footer(slide, 6)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — C2: PPO vs DPO
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "C2: PPO vs. DPO Comparison",
                 "Mean over 3 Seeds | Same Base Model & Reward Model")

# Comparison table
add_table(slide,
    ["Method", "Mean Reward", "Std", "Perplexity", "Alignment Tax", "Efficiency (R/PPL%)"],
    [["SFT (baseline)",
      f"{sft_c2['mean_reward']:.3f}", f"{sft_c2['std_reward']:.3f}",
      f"{sft_c2['perplexity']:.2f}", "—", "—"],
     ["PPO  (β=0.20)",
      f"{ppo_c2['mean_reward']:.3f}", f"{ppo_c2['std_reward']:.3f}",
      f"{ppo_c2['perplexity']:.2f}",
      f"+{method['PPO']['tax_pct']:.1f}%",
      f"{method['PPO']['efficiency_reward_per_ppl_pct']:.4f}"],
     ["DPO  (β=0.10)",
      f"{dpo_c2['mean_reward']:.3f}", f"{dpo_c2['std_reward']:.3f}",
      f"{dpo_c2['perplexity']:.2f}",
      f"+{method['DPO']['tax_pct']:.1f}%",
      f"{method['DPO']['efficiency_reward_per_ppl_pct']:.4f}"]],
    Inches(0.4), Inches(1.4), Inches(7.8), Inches(1.8))

# PPL scatter plot
add_image_safe(slide, FIG_PPL,
               Inches(8.3), Inches(1.3), Inches(4.7), Inches(3.5))
add_text_box(slide, "Fig. 2 — Perplexity vs. Reward per method",
             Inches(8.3), Inches(4.85), Inches(4.7), Inches(0.4),
             font_size=11, italic=True, color=C_GREY, align=PP_ALIGN.CENTER)

# Key insight box
eff_ratio = (method['PPO']['efficiency_reward_per_ppl_pct'] /
             method['DPO']['efficiency_reward_per_ppl_pct'])
add_rect(slide, Inches(0.4), Inches(3.45), Inches(7.8), Inches(1.5),
         fill_rgb=RGBColor(0xE8, 0xF5, 0xE9), line_rgb=C_GREEN, line_width=Pt(1))
add_text_box(slide,
    f"Key Finding:  PPO is {eff_ratio:.1f}× more reward-efficient than DPO",
    Inches(0.6), Inches(3.55), Inches(7.3), Inches(0.55),
    font_size=15, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
add_text_box(slide,
    f"PPO tax: +{method['PPO']['tax_pct']:.1f}%  vs.  DPO tax: +{method['DPO']['tax_pct']:.1f}%  |  "
    f"PPO achieves ≈ same reward ({ppo_c2['mean_reward']:.3f} vs {dpo_c2['mean_reward']:.3f}) "
    f"with much less perplexity increase",
    Inches(0.6), Inches(4.1), Inches(7.3), Inches(0.7),
    font_size=13, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

add_slide_footer(slide, 7)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — C3: Stability Sweep Overview
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "C3: KL-Penalty Stability Sweep",
                 "β ∈ {0.05, 0.10, 0.20, 0.50} × 3 seeds = 12 runs total")

ok    = sweep[sweep["status"] == "ok"]
grp   = ok.groupby("beta").agg(
    mean_r=("mean_reward_after","mean"),
    std_r =("mean_reward_after","std"),
    mean_k=("mean_kl_final_20","mean"),
    max_k =("max_kl","max"),
    coll  =("collapse_detected","sum")).reset_index()

rows = []
for r in grp.itertuples():
    verdict = "⚠ HACKING" if r.coll > 0 else "✓ Stable"
    rows.append([
        f"{r.beta}", f"{r.mean_r:.3f}", f"±{r.std_r:.3f}",
        f"{r.mean_k:.2f}", f"{r.max_k:.2f}", f"{int(r.coll)}/3", verdict
    ])

add_table(slide,
    ["β", "Mean Reward", "±Std", "Mean KL", "Max KL", "Collapse", "Verdict"],
    rows,
    Inches(0.4), Inches(1.45), Inches(7.4), Inches(2.2))

# Phase diagram
add_rect(slide, Inches(8.1), Inches(1.4), Inches(4.8), Inches(2.2),
         fill_rgb=C_LIGHT_BLUE, line_rgb=C_MID_BLUE, line_width=Pt(1))
add_text_box(slide, "Phase Diagram", Inches(8.1), Inches(1.4),
             Inches(4.8), Inches(0.45),
             font_size=15, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

phase_items = [
    (f"β = 0.05 → COLLAPSE  (KL ≈ 43–67)", C_RED),
    (f"β = 0.10 → COLLAPSE  (KL ≈ 23–24)", C_RED),
    (f"β = 0.20 → STABLE ★  (KL ≈ 14.2–14.6)", C_GREEN),
    (f"β = 0.50 → STABLE     (KL ≈ 7.5–8.1)", C_GREEN),
]
for i, (txt, col) in enumerate(phase_items):
    y_p = Inches(1.9) + i * Inches(0.44)
    add_text_box(slide, txt, Inches(8.2), y_p, Inches(4.6), Inches(0.42),
                 font_size=14, bold=(col == C_GREEN and "★" in txt), color=col)

# KL onset box
add_rect(slide, Inches(0.4), Inches(3.85), Inches(12.5), Inches(1.0),
         fill_rgb=RGBColor(0xFF, 0xF0, 0xD0), line_rgb=C_ORANGE, line_width=Pt(1.5))
add_text_box(slide,
    f"Reward Hacking Onset: Mean KL ≥ {opt['collapse_onset_kl']:.1f}  "
    f"|  Safety Margin at β=0.20: {opt['safety_margin_kl']:.1f} KL units  "
    f"|  Optimal β: 0.20",
    Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.7),
    font_size=15, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

add_slide_footer(slide, 8)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Sweep Plots
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "C3: Stability Sweep — Visualisations",
                 "Gao et al. (2023) — proxy reward vs KL divergence")

add_image_safe(slide, FIG_RVK, Inches(0.3), Inches(1.3), Inches(6.2), Inches(4.5))
add_text_box(slide, "Fig. 3 — Reward vs. KL per β (scatter)",
             Inches(0.3), Inches(5.85), Inches(6.2), Inches(0.4),
             font_size=12, italic=True, color=C_GREY, align=PP_ALIGN.CENTER)

add_image_safe(slide, FIG_BETA, Inches(6.8), Inches(1.3), Inches(6.2), Inches(4.5))
add_text_box(slide, "Fig. 4 — Mean reward & KL per β (bar summary)",
             Inches(6.8), Inches(5.85), Inches(6.2), Inches(0.4),
             font_size=12, italic=True, color=C_GREY, align=PP_ALIGN.CENTER)

add_slide_footer(slide, 9)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Alignment Tax Analysis
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "Alignment Tax Analysis",
                 "Perplexity cost of reward alignment — PPO vs DPO")

add_image_safe(slide, FIG_TAX, Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.8))
add_text_box(slide, "Fig. 5 — Alignment tax per β (PPO)",
             Inches(0.3), Inches(5.2), Inches(6.2), Inches(0.4),
             font_size=12, italic=True, color=C_GREY, align=PP_ALIGN.CENTER)

add_image_safe(slide, FIG_OPT, Inches(6.8), Inches(1.3), Inches(6.2), Inches(3.8))
add_text_box(slide, "Fig. 6 — Optimal KL stopping point analysis",
             Inches(6.8), Inches(5.2), Inches(6.2), Inches(0.4),
             font_size=12, italic=True, color=C_GREY, align=PP_ALIGN.CENTER)

# Tax summary
add_rect(slide, Inches(0.3), Inches(5.7), Inches(12.7), Inches(0.9),
         fill_rgb=C_LIGHT_BLUE, line_rgb=C_MID_BLUE, line_width=Pt(1))
add_text_box(slide,
    f"PPO Tax: +{method['PPO']['tax_pct']:.1f}%  |  DPO Tax: +{method['DPO']['tax_pct']:.1f}%  "
    f"|  Optimal Stopping KL: {opt['optimal_mean_kl']:.2f}  "
    f"|  Collapse Onset KL: {opt['collapse_onset_kl']:.2f}  "
    f"|  Safety Margin: {opt['safety_margin_kl']:.2f}",
    Inches(0.4), Inches(5.75), Inches(12.5), Inches(0.75),
    font_size=14, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

add_slide_footer(slide, 10)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Key Findings
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "Key Findings",
                 "Three empirical contributions")

findings = [
    (f"C1 — PPO Reliably Aligns GPT-2",
     f"Mean reward: {agg['sft_mean_reward']:.3f} → {agg['ppo_mean_reward']:.3f}  "
     f"(+{agg['reward_delta']:.3f}, +{agg['alignment_tax_pct']:.1f}% PPL)  |  "
     f"Zero reward collapse across all seeds",
     C_MID_BLUE),
    (f"C2 — PPO Beats DPO in Reward Efficiency",
     f"PPO efficiency: {method['PPO']['efficiency_reward_per_ppl_pct']:.4f}  vs.  "
     f"DPO: {method['DPO']['efficiency_reward_per_ppl_pct']:.4f}  |  "
     f"{method['PPO']['efficiency_reward_per_ppl_pct']/method['DPO']['efficiency_reward_per_ppl_pct']:.1f}× advantage  |  "
     f"PPO's KL constraint provides tighter distribution control",
     C_GREEN),
    (f"C3 — β = 0.20 is the Optimal KL-Penalty",
     f"Hacking onset: KL ≥ {opt['collapse_onset_kl']:.1f}  |  "
     f"Safety margin: {opt['safety_margin_kl']:.1f} units at β=0.20  |  "
     f"Consistent with Gao et al. (2023) scaling laws",
     C_DARK_BLUE),
]

for i, (title, desc, color) in enumerate(findings):
    y_f = Inches(1.4) + i * Inches(1.7)
    add_rect(slide, Inches(0.4), y_f, Inches(12.5), Inches(1.55), fill_rgb=color)
    add_text_box(slide, title, Inches(0.6), y_f + Inches(0.1),
                 Inches(12), Inches(0.6),
                 font_size=18, bold=True, color=C_WHITE)
    add_text_box(slide, desc, Inches(0.6), y_f + Inches(0.7),
                 Inches(12), Inches(0.8),
                 font_size=14, color=C_WHITE)

add_slide_footer(slide, 11)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — Conclusion & Recommendation
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "Conclusion & Recommendation",
                 "Reproducible RLHF pipeline for GPT-2 scale")

bullets = [
    ("Complete PPO-RLHF pipeline on GPT-2/IMDB — all code open-sourced on GitHub", 0),
    (f"Mean reward improvement: +{agg['reward_delta']:.2f} pts (+{agg['alignment_tax_pct']:.1f}% PPL) — stable across 3 seeds", 1),
    ("PPO outperforms DPO at this scale", 0),
    (f"2.5× higher reward-per-perplexity efficiency; lower alignment tax (+{method['PPO']['tax_pct']:.1f}% vs +{method['DPO']['tax_pct']:.1f}%)", 1),
    ("Practical recommendation: use β = 0.20", 0),
    (f"Early-stop rule: halt PPO if mean KL > {opt['collapse_onset_kl']:.1f} (collapse onset threshold)", 1),
    (f"Safety margin: {opt['safety_margin_kl']:.1f} KL units at β=0.20 before hacking onset", 1),
    ("Future Work", 0),
    ("Scale to GPT-2 Medium/Large; true human preference labels; longer training", 1),
    ("Explore GRPO, REINFORCE++, and other RLHF variants", 1),
]
add_bullet_box(slide, bullets, Inches(0.5), Inches(1.45), Inches(12.5), Inches(5.0),
               font_size=16)

# Recommended config box
add_rect(slide, Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.7),
         fill_rgb=C_DARK_BLUE)
add_text_box(slide,
    f"Recommended Config:  model=gpt2  |  β=0.20  |  steps=200  |  "
    f"early_stop_kl={opt['collapse_onset_kl']:.0f}  |  seeds=[42,123,456]",
    Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
    font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_slide_footer(slide, 12)

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — References
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_WHITE)
add_slide_header(slide, "References")

refs = [
    "[1] P. Christiano et al., \"Deep RL from human preferences,\" NeurIPS 2017.",
    "[2] L. Ouyang et al., \"Training language models to follow instructions with RLHF (InstructGPT),\" NeurIPS 2022.",
    "[3] L. Gao et al., \"Scaling laws for reward model overoptimisation,\" ICML 2023.",
    "[4] R. Rafailov et al., \"Direct Preference Optimisation (DPO),\" NeurIPS 2023.",
    "[5] Y. Bai et al., \"Training a helpful and harmless assistant with RLHF,\" arXiv:2204.05862, 2022.",
    "[6] A. Radford et al., \"Language models are unsupervised multitask learners (GPT-2),\" OpenAI Blog, 2019.",
    "[7] A. L. Maas et al., \"Learning word vectors for sentiment analysis (IMDB),\" ACL 2011.",
    "[8] L. von Werra et al., \"TRL: Transformer Reinforcement Learning,\" HuggingFace, 2022.",
    "[9] J. Schulman et al., \"Proximal Policy Optimization Algorithms,\" arXiv:1707.06347, 2017.",
]
ref_lines = [(r, False, 14, C_BLACK) for r in refs]
add_text_box_multi(slide, ref_lines, Inches(0.5), Inches(1.4),
                   Inches(12.5), Inches(5.5),
                   font_size=14, color=C_BLACK)

add_text_box(slide,
    "Muhammad Haseeb  |  Muavia Shakeel  |  FAST University Islamabad  |  Masters in AI",
    Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
    font_size=13, italic=True, color=C_GREY, align=PP_ALIGN.CENTER)

add_slide_footer(slide, 13)

# ═══════════════════════════════════════════════════════════════
out_path = os.path.join(OUT_DIR, "RLHF_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
