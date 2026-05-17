"""
Generates RLHF_Project_Proposal.docx and RLHF_Project_Proposal.pptx
Run: python3 generate_proposal.py
"""

# ─────────────────────────────────────────────────────────────────
# DOCX  ── Word proposal
# ─────────────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import datetime

doc = Document()

# ── Styles helper ──────────────────────────────────────────────
def set_font(run, size=11, bold=False, color=None):
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)

def heading(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    return p

def para(doc, text, bold=False, size=11, align=WD_ALIGN_PARAGRAPH.JUSTIFY):
    p = doc.add_paragraph()
    p.alignment = align
    run = p.add_run(text)
    set_font(run, size, bold)
    return p

def add_bullet(doc, text, style='List Bullet'):
    p = doc.add_paragraph(text, style=style)
    return p

# ── Page margins ───────────────────────────────────────────────
section = doc.sections[0]
section.top_margin    = Inches(1)
section.bottom_margin = Inches(1)
section.left_margin   = Inches(1.25)
section.right_margin  = Inches(1.25)

# ── Title Block ────────────────────────────────────────────────
t = doc.add_paragraph()
t.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = t.add_run("Project Proposal")
set_font(r, 14, True, (31, 73, 125))

t2 = doc.add_paragraph()
t2.alignment = WD_ALIGN_PARAGRAPH.CENTER
r2 = t2.add_run(
    "RLHF for Language Model Alignment:\n"
    "Reproducible GPT-2 Fine-Tuning with PPO,\n"
    "Comparative Analysis, and Stability Benchmarking"
)
set_font(r2, 16, True, (0, 0, 0))

doc.add_paragraph()
info = doc.add_table(rows=4, cols=2)
info.style = 'Table Grid'
labels = ["Student:", "Programme:", "Subject:", "Submission Date:"]
values = ["[Your Name]", "Masters in Artificial Intelligence",
          "Reinforcement Learning", datetime.date.today().strftime("%B %d, %Y")]
for i, (lbl, val) in enumerate(zip(labels, values)):
    info.cell(i, 0).text = lbl
    info.cell(i, 1).text = val
    info.cell(i, 0).paragraphs[0].runs[0].bold = True

doc.add_paragraph()
doc.add_page_break()

# ── 1. Abstract ────────────────────────────────────────────────
heading(doc, "1. Abstract")
para(doc,
    "This project presents a three-contribution study on Reinforcement Learning "
    "from Human Feedback (RLHF) applied to GPT-2 (117 M parameters). "
    "Contribution 1 delivers a complete, reproducible end-to-end RLHF pipeline "
    "(SFT → Reward Model → PPO) using the HuggingFace TRL library on the IMDB "
    "sentiment-control task. Contribution 2 extends this with a head-to-head "
    "comparative analysis of PPO vs. Direct Preference Optimisation (DPO) on "
    "the same task and hardware, quantifying alignment quality versus training "
    "cost. Contribution 3 empirically measures training stability and the "
    "Alignment Tax — the degradation in general language modelling capability "
    "caused by RLHF — by tracking perplexity, KL divergence, and collapse rate "
    "across multiple runs. Future extensions including a behavioral alignment "
    "study (Contribution 4) and active preference learning via ADPO "
    "(Contribution 5) are scoped for subsequent work.")

# ── 2. Problem Statement ───────────────────────────────────────
heading(doc, "2. Problem Statement and Motivation")
para(doc,
    "Large language models (LLMs) trained on raw internet text frequently produce "
    "outputs that are factually incorrect, unhelpful, or misaligned with human "
    "intent. OpenAI's InstructGPT (Ouyang et al., 2022) demonstrated that a "
    "relatively small 1.3 B-parameter model fine-tuned with RLHF can outperform "
    "a 175 B-parameter model trained without RLHF on the majority of human "
    "preference tasks. This result established RLHF as the de-facto alignment "
    "training paradigm and underpins the training of ChatGPT, Claude, and Gemini.")
para(doc,
    "Despite its industrial impact, most RLHF demonstrations use proprietary "
    "large-scale models, making them inaccessible for academic study. Furthermore, "
    "the comparative behaviour of PPO vs. newer preference-optimisation methods "
    "at small scale remains understudied, and the practical stability cost "
    "('Alignment Tax') of applying RLHF is rarely measured in open, reproducible "
    "settings. This project addresses all three gaps at GPT-2 scale using fully "
    "open-source tools.")

# ── 3. Research Questions ──────────────────────────────────────
heading(doc, "3. Research Questions")
para(doc, "This project addresses three inter-related research questions:", bold=False)
for rq in [
    "RQ1: How does RLHF with PPO improve GPT-2's alignment on a sentiment-control "
    "task, and what is the measurable improvement in reward score and output quality "
    "compared to a supervised fine-tuning baseline?",
    "RQ2: Does Direct Preference Optimisation (DPO) achieve comparable alignment "
    "to PPO on the same task at lower computational cost, and what are the "
    "quantitative trade-offs in reward score, perplexity, and training time?",
    "RQ3: What is the empirical Alignment Tax of RLHF at GPT-2 scale — how much "
    "general language modelling capability (perplexity) is sacrificed for alignment "
    "gain, and at what KL divergence does reward hacking onset occur?",
]:
    add_bullet(doc, rq)

# ── 4. Background ─────────────────────────────────────────────
heading(doc, "4. Background and Related Work")
para(doc,
    "The following key works ground this project:")

refs_summary = [
    ("Stiennon et al. (2020)", "Introduced learning to summarise with human feedback, "
     "establishing the RLHF three-stage pipeline for NLP."),
    ("Ouyang et al. (2022) — InstructGPT", "Showed RLHF-aligned 1.3 B model outperforms "
     "GPT-3 175 B on human preference; the canonical RLHF reference."),
    ("Bai et al. (2022) — HH Dataset", "Released the Helpful & Harmless dataset; "
     "demonstrated Constitutional AI and RLHF at scale."),
    ("Schulman et al. (2017) — PPO", "Introduced Proximal Policy Optimisation, the "
     "clipped-objective actor-critic algorithm used in RLHF policy training."),
    ("Gao et al. (2023)", "Studied reward model over-optimisation ('reward hacking'), "
     "showing KL-regularisation is critical and defining the scaling law for optimal "
     "stopping — directly motivating Contribution 3."),
    ("Zheng et al. (2023) — Secrets of RLHF", "Identified the engineering decisions "
     "(reward normalisation, value clipping, adaptive KL) that determine PPO stability "
     "— the foundation for Contribution 3's stability benchmarking."),
    ("Rafailov et al. (2023) — DPO", "Reformulated RLHF as direct supervised learning "
     "over preference pairs, eliminating the RL loop — the method compared in "
     "Contribution 2. Achieves comparable alignment at 10× lower training cost."),
    ("Touvron et al. (2023) — Llama 2", "Released full RLHF training details including "
     "PPOConfig hyper-parameters adopted in this project; motivates dual-RM as future work."),
    ("Zhu et al. (2023) — APA", "Demonstrated PPO collapse rate of 23% on GPT-2 scale "
     "sentiment control — provides direct baselines for Contributions 1 and 2."),
]
for name, summary in refs_summary:
    p = doc.add_paragraph(style='List Bullet')
    r_bold = p.add_run(name + ": ")
    r_bold.bold = True
    r_bold.font.size = Pt(11)
    r_normal = p.add_run(summary)
    r_normal.font.size = Pt(11)

# ── 5. Contributions ──────────────────────────────────────────
heading(doc, "5. Project Contributions")

heading(doc, "5.1  Contribution 1: Reproducible GPT-2 RLHF Pipeline", level=2)
para(doc,
    "We implement a complete, documented, reproducible RLHF pipeline at GPT-2 "
    "scale using the HuggingFace TRL library. The three-stage pipeline covers:")
for b in [
    "Stage 1 — Supervised Fine-Tuning (SFT): Fine-tune GPT-2 on positive IMDB reviews.",
    "Stage 2 — Reward Model Training (RM): Train a GPT-2-based sequence classifier "
    "to score positive vs. negative sentiment as a proxy reward signal.",
    "Stage 3 — PPO Policy Optimisation: Use TRL PPOTrainer to optimise GPT-2 against "
    "the RM reward signal with KL-divergence penalty (β = 0.2).",
    "Deliverables: Mean reward score before/after PPO; reward learning curve (200 steps); "
    "KL divergence monitoring; 10+ qualitative prompt–response pairs; interactive CLI demo.",
]:
    add_bullet(doc, b)

heading(doc, "5.2  Contribution 2: PPO vs. DPO Comparative Analysis", level=2)
para(doc,
    "Using the pipeline from Contribution 1 as the PPO baseline, we train an "
    "equivalent GPT-2 model with Direct Preference Optimisation (DPO; Rafailov "
    "et al., 2023) on the same IMDB sentiment-control task and identical hardware. "
    "This is a true head-to-head comparison — same dataset, same base model, same "
    "evaluation protocol. We measure and compare:")
for b in [
    "Reward score (mean ± std) before and after alignment for both methods",
    "Perplexity on held-out IMDB reviews (Alignment Tax indicator)",
    "Wall-clock training time (GPU hours)",
    "Training collapse rate across 3 seeds per method",
    "Qualitative response quality: 10 prompts evaluated under SFT / PPO / DPO",
]:
    add_bullet(doc, b)
para(doc,
    "Expected outcome: DPO achieves comparable reward score to PPO with "
    "approximately 10× faster training, consistent with Rafailov et al. (2023).")

heading(doc, "5.3  Contribution 3: Training Stability and Alignment Tax Measurement", level=2)
para(doc,
    "We systematically benchmark the stability and capability cost of RLHF at "
    "GPT-2 scale. This contribution provides empirical data on the 'Alignment Tax' "
    "— the degradation in general language modelling performance caused by reward "
    "optimisation. Specifically:")
for b in [
    "Run 3 seeds per method (PPO, DPO) to measure collapse rate and variance",
    "Vary KL coefficient β ∈ {0.05, 0.1, 0.2, 0.5} for PPO to identify the "
    "reward hacking onset point (per Gao et al., 2023 scaling law)",
    "Track perplexity increase vs. reward score gain across training steps "
    "(the Alignment Tax curve)",
    "Identify the optimal stopping point: the KL threshold at which gold-metric "
    "improvement plateaus before reward hacking begins",
    "Report: recommended stable training configuration for GPT-2-scale RLHF",
]:
    add_bullet(doc, b)
para(doc,
    "Expected outcome: PPO collapse rate of ~20–25% (baseline from Zhu et al., 2023) "
    "reduced to <5% with optimal β and reward normalisation. Alignment Tax of ≤15% "
    "perplexity increase at the recommended stopping point.")

# ── 6. Methodology ────────────────────────────────────────────
heading(doc, "6. Methodology")

heading(doc, "6.1  Dataset", level=2)
para(doc,
    "IMDB Movie Reviews dataset (50,000 reviews, HuggingFace datasets). "
    "Positive reviews (label=1) serve as the SFT corpus. Positive/negative pairs "
    "serve as the reward model training data and DPO preference pairs. "
    "Prompts for PPO/DPO evaluation are sampled from the first 16 tokens of "
    "test-split reviews.")

heading(doc, "6.2  Models and Libraries", level=2)
for item in [
    "Base model: GPT-2 (124 M parameters, HuggingFace hub: gpt2)",
    "Reward model: GPT-2 + sequence-classification head (Contributions 1 & 3)",
    "PPO Policy: GPT-2 initialised from SFT checkpoint + TRL PPOTrainer",
    "DPO Policy: GPT-2 initialised from SFT checkpoint + TRL DPOTrainer",
    "Reference model: Frozen SFT checkpoint (KL target for both PPO and DPO)",
    "Library: HuggingFace TRL, transformers, datasets, torch",
]:
    add_bullet(doc, item)

heading(doc, "6.3  PPO Hyper-parameters", level=2)
for item in [
    "Learning rate: 1.41 × 10⁻⁵  (from Llama 2 paper, Touvron et al., 2023)",
    "Batch size: 16 | Mini-batch size: 4 | PPO epochs: 4",
    "Max new tokens: 32 | Training steps: 200",
    "KL coefficient β: 0.2 (swept over {0.05, 0.1, 0.2, 0.5} for Contribution 3)",
]:
    add_bullet(doc, item)

heading(doc, "6.4  DPO Hyper-parameters", level=2)
for item in [
    "Learning rate: 5 × 10⁻⁵",
    "β (temperature): 0.1 (DPO reference divergence parameter)",
    "Batch size: 16 | Training epochs: 3",
    "Loss type: sigmoid (standard DPO)",
]:
    add_bullet(doc, item)

heading(doc, "6.5  Evaluation Metrics", level=2)
for item in [
    "Mean reward score ± std (before vs. after alignment, per method)",
    "Reward standard deviation across seeds (stability proxy)",
    "KL divergence from SFT reference (policy drift, per training step)",
    "Perplexity on held-out IMDB reviews (Alignment Tax)",
    "Training collapse rate (% of runs with reward → 0 or perplexity > 200)",
    "Wall-clock training time (GPU hours) per method",
    "Qualitative response comparison (10 prompts × SFT / PPO / DPO)",
]:
    add_bullet(doc, item)

# ── 7. Timeline ───────────────────────────────────────────────
heading(doc, "7. Ten-Day Implementation Timeline")
days = [
    ("Day 1", "Environment Setup & SFT",
     "Install TRL / HuggingFace stack; prepare IMDB dataset; "
     "fine-tune GPT-2 SFT checkpoint; validate loss convergence."),
    ("Day 2", "Reward Model Training",
     "Train reward model (GPT-2 + cls head) on IMDB pos/neg pairs; "
     "evaluate RM accuracy on held-out set; save checkpoint."),
    ("Day 3", "PPO Training — Contribution 1",
     "Configure PPOTrainer with β=0.2; run 200-step PPO (3 seeds); "
     "log reward curve, KL divergence, collapse rate."),
    ("Day 4", "PPO Evaluation — Contribution 1",
     "Compute reward score before/after; plot reward vs. step curve; "
     "generate 10+ before/after qualitative pairs; build CLI demo."),
    ("Day 5", "DPO Training — Contribution 2",
     "Prepare preference pairs from IMDB; configure DPOTrainer; "
     "run DPO training (3 seeds); log equivalent metrics."),
    ("Day 6", "PPO vs DPO Comparison — Contribution 2",
     "Build side-by-side comparison table: reward score, perplexity, "
     "training time, collapse rate. Generate qualitative SFT/PPO/DPO table."),
    ("Day 7", "Stability Sweep — Contribution 3",
     "Run PPO with β ∈ {0.05, 0.1, 0.2, 0.5}; "
     "identify reward hacking onset per Gao et al. (2023) scaling law."),
    ("Day 8", "Alignment Tax Analysis — Contribution 3",
     "Compute perplexity vs. reward gain curves; identify optimal stopping KL; "
     "summarise recommended stable training configuration."),
    ("Day 9", "Write-Up: Methodology + Results",
     "Draft Methodology, Results, and Discussion sections "
     "with all tables and figures populated."),
    ("Day 10", "Demo, Documentation & Submission",
     "Polish interactive demo; write README; "
     "finalise report and appendices; submit code + report."),
]
tbl = doc.add_table(rows=1, cols=3)
tbl.style = 'Table Grid'
for i, hdr in enumerate(["Day", "Milestone", "Tasks"]):
    cell = tbl.cell(0, i)
    cell.text = hdr
    cell.paragraphs[0].runs[0].bold = True
for day, milestone, tasks in days:
    row = tbl.add_row()
    row.cells[0].text = day
    row.cells[1].text = milestone
    row.cells[2].text = tasks

doc.add_paragraph()

# ── 8. Future Work ────────────────────────────────────────────
heading(doc, "8. Future Work")
para(doc,
    "The following two contributions are explicitly planned as extensions in "
    "subsequent research beyond this project submission:")

heading(doc, "8.1  Contribution 4: Investigating 'Human-Likeness' Post-Alignment", level=2)
para(doc,
    "A 2025 finding by Binz et al. demonstrates that post-training (RLHF/instruction "
    "tuning) consistently makes models less human-like in their behavioral patterns, "
    "particularly in reasoning tasks. This contribution would evaluate the model not "
    "only on reward score but on whether its reasoning patterns become more or less "
    "similar to human behavioral transcripts, using the Psych-201 dataset. "
    "This provides a richer 'behavioral alignment' assessment beyond scalar reward "
    "scores, connecting alignment research to cognitive science.")

heading(doc, "8.2  Contribution 5: Efficient Alignment via Active Learning (ADPO)", level=2)
para(doc,
    "Most RLHF projects assume a fixed dataset of preference labels. Active Direct "
    "Preference Optimisation (ADPO) proposes selectively querying the reward model "
    "only when it is uncertain about a response pair, demonstrating that models can "
    "reach the same performance level using significantly fewer labels (estimated "
    "40–60% label reduction). This contribution would implement a simplified ADPO "
    "loop on GPT-2, validating label-efficiency claims at small scale and providing "
    "practical guidance for resource-constrained alignment research.")

# ── 9. Expected Results ───────────────────────────────────────
heading(doc, "9. Expected Results")
para(doc,
    "Based on prior work (Zheng et al., 2023; Gao et al., 2023; Rafailov et al., 2023), "
    "we anticipate the following outcomes:")

heading(doc, "9.1  Contribution 1 — PPO Pipeline", level=2)
for item in [
    "Mean reward score improvement of 40–70% after 200 PPO steps.",
    "Stable KL divergence below 10 nats throughout training with β = 0.2.",
    "Qualitative improvement: responses shift toward positive sentiment.",
    "Full pipeline executable on a single GPU in under 2 hours.",
]:
    add_bullet(doc, item)

heading(doc, "9.2  Contribution 2 — PPO vs. DPO", level=2)
for item in [
    "DPO achieves comparable reward score to PPO (within ±5%).",
    "DPO trains approximately 10× faster than PPO.",
    "DPO has lower collapse rate (~2%) vs. PPO (~20%) across seeds.",
    "PPO achieves higher peak reward score due to online exploration.",
]:
    add_bullet(doc, item)

heading(doc, "9.3  Contribution 3 — Stability & Alignment Tax", level=2)
for item in [
    "PPO collapse rate reduced from ~23% to <5% with reward normalisation and β=0.2.",
    "Optimal stopping point: KL ≈ 10–15 nats before reward hacking onset.",
    "Alignment Tax: ≤15% perplexity increase at the recommended stopping point.",
    "β ∈ {0.05, 0.1}: faster reward gain but earlier reward hacking.",
    "β = 0.5: high stability but slow reward improvement; diminishing alignment.",
]:
    add_bullet(doc, item)

# ── 10. References ────────────────────────────────────────────
heading(doc, "10. References")
references = [
    "Ziegler, D., Stiennon, N., Wu, J., Brown, T., Radford, A., Amodei, D., Christiano, P., & Irving, G. (2019). Fine-tuning language models from human preferences. arXiv. https://arxiv.org/abs/1909.08593",
    "Stiennon, N., Ouyang, L., Wu, J., Ziegler, D., Lowe, R., Voss, C., Radford, A., Amodei, D., & Christiano, P. (2020). Learning to summarize with human feedback. Advances in Neural Information Processing Systems, 33, 3008\u20133021.",
    "Christiano, P., Leike, J., Brown, T., Martic, M., Legg, S., & Amodei, D. (2017). Deep reinforcement learning from human preferences. Advances in Neural Information Processing Systems, 30.",
    "Schulman, J., Wolski, F., Dhariwal, P., Radford, A., & Klimov, O. (2017). Proximal policy optimization algorithms. arXiv. https://arxiv.org/abs/1707.06347",
    "Brown, T., Mann, B., Ryder, N., Subbiah, M., Kaplan, J., Dhariwal, P., Neelakantan, A., Shyam, P., Sastry, G., Askell, A., Agarwal, S., Herbert-Voss, A., Krueger, G., Henighan, T., Child, R., Ramesh, A., Ziegler, D., Wu, J., Winter, C., \u2026 Amodei, D. (2020). Language models are few-shot learners. Advances in Neural Information Processing Systems, 33, 1877\u20131901.",
    "Ouyang, L., Wu, J., Jiang, X., Almeida, D., Wainwright, C., Mishkin, P., Zhang, C., Agarwal, S., Slama, K., Ray, A., Schulman, J., Hilton, J., Kelton, F., Miller, L., Simens, M., Askell, A., Welinder, P., Christiano, P., Leike, J., & Lowe, R. (2022). Training language models to follow instructions with human feedback. Advances in Neural Information Processing Systems, 35, 27730\u201327744.",
    "Bai, Y., Jones, A., Ndousse, K., Askell, A., Chen, A., DasSarma, N., Drain, D., Fort, S., Ganguli, D., Henighan, T., Joseph, N., Kadavath, S., Kernion, J., Conerly, T., El-Showk, S., Elhage, N., Hatfield-Dodds, Z., Hernandez, D., Hume, T., \u2026 Kaplan, J. (2022). Training a helpful and harmless assistant with reinforcement learning from human feedback. arXiv. https://arxiv.org/abs/2204.05862",
    "Gao, L., Biderman, S., Black, S., Golding, L., Hoppe, T., Foster, C., Phang, J., He, H., Thite, A., Nabeshima, N., Presser, S., & Leahy, C. (2023). Scaling laws for reward model overoptimization. Proceedings of the 40th International Conference on Machine Learning.",
    "Rafailov, R., Sharma, A., Mitchell, E., Manning, C. D., Ermon, S., & Finn, C. (2023). Direct preference optimization: Your language model is secretly a reward model. Advances in Neural Information Processing Systems, 36.",
    "Touvron, H., Martin, L., Stone, K., Albert, P., Almahairi, A., Babaei, Y., Bashlykov, N., Batra, S., Bhargava, P., Bhosale, S., Bikel, D., Blecher, L., Ferrer, C. C., Chen, M., Cucurull, G., Esiobu, D., Fernandes, J., Fu, J., Fu, W., \u2026 Scialom, T. (2023). Llama 2: Open foundation and fine-tuned chat models. arXiv. https://arxiv.org/abs/2307.09288",
    "Zheng, R., Dou, S., Gao, S., Hua, Y., Shen, W., Wang, B., Liu, Y., Jin, S., Liu, Q., Zhou, Y., Xiong, L., Chen, L., Xi, Z., Liu, N., Yan, H., Deng, X., Chen, Q., Zhou, E., Gui, T., \u2026 Huang, X. (2023). Secrets of RLHF in large language models part I: PPO. arXiv. https://arxiv.org/abs/2307.04964",
    "Yuan, Z., Yuan, H., Tan, C., Wang, W., Lian, S., & Yu, T. (2023). RRHF: Rank responses to align language models with human feedback. Advances in Neural Information Processing Systems, 36.",
    "Dong, H., Xiong, W., Goyal, D., Pan, R., Diao, S., Zhang, J., Shum, K., & Zhang, T. (2023). RAFT: Reward ranked finetuning for generative foundation model alignment. Transactions on Machine Learning Research.",
    "Azar, M. G., Guo, Z. D., Piot, B., Munos, R., Rowland, M., Valko, M., & Calandriello, D. (2024). A general theoretical paradigm to understand learning from human feedback. Proceedings of the 27th International Conference on Artificial Intelligence and Statistics.",
    "Wang, Z., Bi, B., Pentyala, S., Ramnath, K., Du, S., Shrivastava, D., Jiang, X., Lan, Z., & Natarajan, P. (2024). A comprehensive survey of LLM alignment techniques: RLHF, RLAIF, PPO, DPO and more. arXiv. https://arxiv.org/abs/2407.16216",
    "Binz, M., Dasgupta, I., Jagadish, A., Botvinick, M., Fleming, S. M., & Schulz, E. (2025). Turning large language models into cognitive models. arXiv. https://arxiv.org/abs/2306.03917",
    "Liu, Z., Ji, H., Li, H., & Xiong, C. (2024). Active direct preference optimization. arXiv. https://arxiv.org/abs/2402.10141",
]
for ref in references:
    p = doc.add_paragraph(style='List Bullet')
    r = p.add_run(ref)
    r.font.size = Pt(10)

# ── Save DOCX ─────────────────────────────────────────────────
doc_path = "RLHF_Project_Proposal.docx"
doc.save(doc_path)
print(f"✅  Saved: {doc_path}")


# ─────────────────────────────────────────────────────────────────
# PPTX  ── PowerPoint slides  (professional academic redesign)
# ─────────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Professional academic colour palette ──────────────────────
NAVY       = RGBColor(0x1A, 0x2E, 0x4A)   # #1A2E4A  dark navy
CHARCOAL   = RGBColor(0x2C, 0x3E, 0x50)   # #2C3E50  charcoal
STEEL_BLUE = RGBColor(0x29, 0x80, 0xB9)   # #2980B9  steel blue
TEAL       = RGBColor(0x16, 0x78, 0x8A)   # #16788A  teal accent
GOLD       = RGBColor(0xC9, 0xA8, 0x4C)   # #C9A84C  gold highlight
SILVER     = RGBColor(0xEC, 0xF0, 0xF1)   # #ECF0F1  light silver
WARM_WHITE = RGBColor(0xF8, 0xF9, 0xFA)   # #F8F9FA  warm white
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MID_GRAY   = RGBColor(0x95, 0xA5, 0xA6)   # #95A5A6  medium gray
DARK_GRAY  = RGBColor(0x2C, 0x3E, 0x50)   # body text
RULE_GRAY  = RGBColor(0xD5, 0xDC, 0xE1)   # thin rule lines
GREEN      = RGBColor(0x1A, 0x7A, 0x47)   # #1A7A47  success green
RED        = RGBColor(0xA9, 0x3D, 0x3D)   # #A93D3D  warning red

blank_layout = prs.slide_layouts[6]   # completely blank

SLIDE_W = 13.33
SLIDE_H = 7.5
HEADER_H = 0.95   # header bar height
FOOTER_H = 0.35   # footer bar height

# ── Drawing primitives ────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line=None, line_pt=0.75):
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape

def tb(slide, text, l, t, w, h,
       sz=13, bold=False, italic=False, color=DARK_GRAY,
       align=PP_ALIGN.LEFT, wrap=True):
    from pptx.util import Inches, Pt
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(sz)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def bullets(slide, lines, l, t, w, h,
            sz=12.5, color=DARK_GRAY, indent=True,
            line_space_pt=5, first_bold=False):
    """Add a text box with bullet lines (dot prefix). Returns textbox."""
    from pptx.util import Inches, Pt
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(line_space_pt)
        if indent:
            p.level = 0
        run = p.add_run()
        prefix = "\u2022  " if indent else ""
        run.text = prefix + line
        run.font.size = Pt(sz)
        run.font.color.rgb = color
        run.font.bold = (first_bold and i == 0)
    return txb

def add_slide_chrome(slide, title_text, slide_num, subtitle=None):
    """Standard header + footer for content slides."""
    # Header bar
    rect(slide, 0, 0, SLIDE_W, HEADER_H, fill=NAVY)
    # Gold accent rule under header
    rect(slide, 0, HEADER_H, SLIDE_W, 0.04, fill=GOLD)
    # Footer bar
    rect(slide, 0, SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H, fill=SILVER)
    # Title text
    tb(slide, title_text,
       0.45, 0.12, 10.5, HEADER_H - 0.15,
       sz=22, bold=True, color=WHITE)
    # Slide number (right side of header)
    tb(slide, str(slide_num),
       12.7, 0.2, 0.55, 0.55,
       sz=11, color=GOLD, align=PP_ALIGN.RIGHT)
    # Footer caption
    tb(slide, "RLHF for Language Model Alignment  |  Masters in AI  |  2025",
       0.4, SLIDE_H - FOOTER_H + 0.05, 12.5, FOOTER_H - 0.1,
       sz=8.5, color=CHARCOAL, italic=True)

def section_label(slide, text, l, t, w=3.0):
    """Small steel-blue section label (eyebrow text)."""
    tb(slide, text.upper(), l, t, w, 0.3,
       sz=8.5, bold=True, color=STEEL_BLUE)

def thin_rule(slide, l, t, w, color=RULE_GRAY):
    rect(slide, l, t, w, 0.018, fill=color)

def left_accent(slide, l, t, h, color=STEEL_BLUE):
    """Left vertical accent bar for contribution cards."""
    rect(slide, l, t, 0.06, h, fill=color)

# ── TABLE HELPER ─────────────────────────────────────────────
def add_table(slide, rows_data, col_widths, l, t, header_fill=NAVY,
              row_alt=WARM_WHITE, header_text_sz=11, body_text_sz=10.5):
    """
    rows_data: list of lists of strings.
    First row is treated as header.
    col_widths: list of floats (inches).
    """
    from pptx.util import Inches, Pt
    nrows = len(rows_data)
    ncols = len(rows_data[0])
    total_w = sum(col_widths)
    tbl = slide.shapes.add_table(
        nrows, ncols,
        Inches(l), Inches(t),
        Inches(total_w), Inches(nrows * 0.38)).table
    # set col widths
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(header_text_sz if ri == 0 else body_text_sz)
            run.font.bold = (ri == 0)
            run.font.color.rgb = WHITE if ri == 0 else DARK_GRAY
            fill = cell.fill
            fill.solid()
            if ri == 0:
                fill.fore_color.rgb = header_fill
            elif ri % 2 == 1:
                fill.fore_color.rgb = WARM_WHITE
            else:
                fill.fore_color.rgb = SILVER
    return tbl


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
# Full dark background
rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
# Top gold rule
rect(slide, 0, 0, SLIDE_W, 0.06, fill=GOLD)
# Bottom gold rule
rect(slide, 0, SLIDE_H - 0.06, SLIDE_W, 0.06, fill=GOLD)
# Subtle silver panel behind content
rect(slide, 0.6, 1.8, 12.1, 4.2, fill=CHARCOAL)
rect(slide, 0.6, 1.8, 0.06, 4.2, fill=GOLD)   # left gold bar

tb(slide,
   "Masters in Artificial Intelligence  \u2022  Reinforcement Learning  \u2022  2025",
   0.8, 1.2, 11.7, 0.5,
   sz=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

tb(slide,
   "RLHF for Language Model Alignment",
   0.8, 2.0, 11.7, 1.1,
   sz=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tb(slide,
   "Reproducible GPT-2 Fine-Tuning with PPO,\nComparative Analysis & Stability Benchmarking",
   0.8, 3.15, 11.7, 1.2,
   sz=19, italic=True, color=SILVER, align=PP_ALIGN.CENTER)

thin_rule(slide, 1.5, 4.55, 10.3, color=GOLD)

tb(slide,
   "3 Main Contributions:   C1 \u2014 Reproducible Pipeline     C2 \u2014 PPO vs DPO     C3 \u2014 Stability Benchmarking",
   0.8, 4.68, 11.7, 0.55,
   sz=13, color=SILVER, align=PP_ALIGN.CENTER)

tb(slide, "[Student Name]  |  [University Name]",
   0.8, 6.7, 11.7, 0.5,
   sz=11, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Motivation
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "Motivation: Why RLHF?", 2)
BODY_TOP = HEADER_H + 0.12

bullets(slide, [
    "LLMs trained on raw internet text produce harmful, incorrect, or misaligned outputs",
    "RLHF is the alignment paradigm behind ChatGPT, Claude, and Gemini",
    "InstructGPT (Ouyang et al., 2022): a 1.3B RLHF model outperforms GPT-3 175B on human preference",
    "Core shift: reward human preferences, not just next-token prediction",
    "Challenge: most RLHF implementations require proprietary large-scale infrastructure",
    "This project: reproduce the full RLHF loop at GPT-2 scale — open, transparent, and academically rigorous",
], 0.5, BODY_TOP, 8.0, 5.5, sz=14)

# Quote card
rect(slide, 8.7, BODY_TOP, 4.3, 1.9, fill=WARM_WHITE, line=RULE_GRAY)
left_accent(slide, 8.7, BODY_TOP, 1.9, color=GOLD)
tb(slide,
   '"A 1.3B RLHF model is preferred\nover a 175B GPT-3 model by\nhuman evaluators."',
   8.85, BODY_TOP + 0.12, 4.1, 1.4,
   sz=12, italic=True, color=CHARCOAL)
tb(slide, "\u2014 Ouyang et al., 2022 (InstructGPT)",
   8.85, BODY_TOP + 1.6, 4.1, 0.35,
   sz=9.5, italic=True, color=MID_GRAY)

# Pipeline diagram (text-based)
rect(slide, 8.7, BODY_TOP + 2.2, 4.3, 2.6, fill=WARM_WHITE, line=RULE_GRAY)
left_accent(slide, 8.7, BODY_TOP + 2.2, 2.6, color=STEEL_BLUE)
tb(slide, "3-Stage RLHF Pipeline", 8.85, BODY_TOP + 2.28, 4.1, 0.4, sz=11, bold=True, color=NAVY)
for i, step in enumerate(["Stage 1  \u2014  Supervised Fine-Tuning (SFT)",
                           "Stage 2  \u2014  Reward Model Training (RM)",
                           "Stage 3  \u2014  PPO Policy Optimisation"]):
    tb(slide, step, 8.85, BODY_TOP + 2.75 + i * 0.6, 4.1, 0.55, sz=11, color=CHARCOAL)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement & Research Questions
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "Problem Statement & Research Questions", 3)
BODY_TOP = HEADER_H + 0.14

# Problem box
rect(slide, 0.45, BODY_TOP, 12.4, 1.0, fill=WARM_WHITE, line=RULE_GRAY)
left_accent(slide, 0.45, BODY_TOP, 1.0, color=GOLD)
section_label(slide, "Problem", 0.6, BODY_TOP + 0.05)
tb(slide,
   "Most RLHF research operates at scales inaccessible to academic settings. "
   "No compact, reproducible baseline exists comparing PPO vs. DPO at GPT-2 scale, "
   "nor measuring the practical Alignment Tax in an open, repeatable setting.",
   0.6, BODY_TOP + 0.25, 12.1, 0.7, sz=13, color=CHARCOAL)

# RQs
rqs = [
    ("RQ1",
     "How does RLHF with PPO improve GPT-2 alignment on a sentiment-control task, "
     "and what is the measurable improvement in reward score compared to an SFT baseline?"),
    ("RQ2",
     "Does Direct Preference Optimisation (DPO) achieve comparable alignment to PPO "
     "on the same task at lower computational cost, and what are the quantitative trade-offs?"),
    ("RQ3",
     "What is the empirical Alignment Tax at GPT-2 scale \u2014 how much general language "
     "modelling capability is lost, and at what KL divergence does reward hacking onset occur?"),
]
for i, (label, text) in enumerate(rqs):
    ty = BODY_TOP + 1.2 + i * 1.6
    rect(slide, 0.45, ty, 12.4, 1.4, fill=WARM_WHITE, line=RULE_GRAY)
    left_accent(slide, 0.45, ty, 1.4, color=STEEL_BLUE)
    # RQ badge
    rect(slide, 0.55, ty + 0.35, 0.7, 0.55, fill=STEEL_BLUE)
    tb(slide, label, 0.56, ty + 0.37, 0.68, 0.5, sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, text, 1.35, ty + 0.2, 11.3, 1.0, sz=12.5, color=CHARCOAL)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Key Literature
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "Key Literature", 4)
BODY_TOP = HEADER_H + 0.1

papers = [
    ("Stiennon et al. (2020)", "Introduced the RLHF 3-stage NLP pipeline; canonical foundation."),
    ("Ouyang et al. (2022)",   "InstructGPT: RLHF for instruction following; basis of ChatGPT."),
    ("Christiano et al. (2017)", "Deep RL from human preferences; seminal RLHF framework."),
    ("Schulman et al. (2017)", "PPO algorithm: clipped-objective actor-critic used in our training."),
    ("Gao et al. (2023)",      "Reward model overoptimisation: scaling laws motivate KL control."),
    ("Zheng et al. (2023)",    "Secrets of RLHF: open-source PPO stability insights at GPT scale."),
    ("Rafailov et al. (2023)", "DPO: eliminates RL loop; baseline algorithm compared in C2."),
    ("Bai et al. (2022)",      "Helpful & Harmless dataset; Constitutional AI RLHF at scale."),
    ("Brown et al. (2020)",    "GPT-3 few-shot learning; provides scaling context for our work."),
    ("Wang et al. (2024)",     "Comprehensive survey: RLHF, RLAIF, PPO, DPO and extensions."),
]
# Two columns, 5 rows each
for idx, (title, desc) in enumerate(papers):
    col = idx % 2
    row = idx // 2
    lx = 0.45 + col * 6.45
    ty = BODY_TOP + row * 1.18
    rect(slide, lx, ty, 6.15, 1.05, fill=WARM_WHITE, line=RULE_GRAY)
    left_accent(slide, lx, ty, 1.05, color=STEEL_BLUE)
    tb(slide, title,  lx + 0.13, ty + 0.06, 5.9, 0.38, sz=11.5, bold=True, color=NAVY)
    tb(slide, desc,   lx + 0.13, ty + 0.48, 5.9, 0.52, sz=10.5, color=CHARCOAL)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Contribution 1: Reproducible RLHF Pipeline
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "Contribution 1 \u2014 Reproducible GPT-2 RLHF Pipeline", 5)
BODY_TOP = HEADER_H + 0.14

section_label(slide, "C1 Overview", 0.5, BODY_TOP)

# Pipeline stages
stages = [
    ("Stage 1\nSFT",          "Fine-tune GPT-2 on\npositive IMDB reviews"),
    ("Stage 2\nReward Model",  "GPT-2 classifier:\npositive vs negative\nsentiment scoring"),
    ("Stage 3\nPPO",           "TRL PPOTrainer:\nKL-regularised policy\noptimisation (beta=0.2)"),
]
for i, (label, desc) in enumerate(stages):
    lx = 0.45 + i * 3.65
    rect(slide, lx, BODY_TOP + 0.32, 3.25, 2.0, fill=NAVY)
    tb(slide, label,  lx + 0.15, BODY_TOP + 0.38, 2.95, 0.75, sz=13.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, desc,   lx + 0.15, BODY_TOP + 1.12, 2.95, 1.0,  sz=11,   color=SILVER,align=PP_ALIGN.CENTER)
    if i < 2:
        tb(slide, "\u25b6", 3.55 + i * 3.65, BODY_TOP + 0.85, 0.35, 0.5, sz=14, bold=True, color=GOLD)

# Divider
thin_rule(slide, 0.45, BODY_TOP + 2.55, 12.4)

# Left: Deliverables
section_label(slide, "Deliverables", 0.5, BODY_TOP + 2.65)
bullets(slide, [
    "Mean reward score before vs. after PPO (200 training steps)",
    "Reward learning curve and KL divergence trajectory",
    "Perplexity comparison: SFT baseline vs. PPO-aligned model",
    "Qualitative prompt-response pairs (>10 examples)",
    "Interactive CLI demo for live inference",
], 0.5, BODY_TOP + 3.0, 6.5, 2.9, sz=12.5)

# Right: Expected Demo
rect(slide, 7.3, BODY_TOP + 2.65, 5.6, 3.25, fill=WARM_WHITE, line=RULE_GRAY)
section_label(slide, "Expected Output Example", 7.45, BODY_TOP + 2.72)
tb(slide, 'Prompt:  "This movie was..."', 7.45, BODY_TOP + 3.0, 5.3, 0.38, sz=11.5, bold=True, color=CHARCOAL)
rect(slide, 7.45, BODY_TOP + 3.45, 5.3, 0.65, fill=RGBColor(0xF9, 0xEB, 0xEB), line=RULE_GRAY)
left_accent(slide, 7.45, BODY_TOP + 3.45, 0.65, color=RED)
tb(slide, "Before PPO (SFT):   reward \u2248 \u22120.8\n\"...slow, tedious, and entirely forgettable.\"",
   7.6, BODY_TOP + 3.5, 5.05, 0.55, sz=10.5, color=DARK_GRAY)
rect(slide, 7.45, BODY_TOP + 4.2, 5.3, 0.65, fill=RGBColor(0xEA, 0xF6, 0xEC), line=RULE_GRAY)
left_accent(slide, 7.45, BODY_TOP + 4.2, 0.65, color=GREEN)
tb(slide, "After PPO (aligned):   reward \u2248 +2.1\n\"...a masterpiece. Absolutely brilliant storytelling!\"",
   7.6, BODY_TOP + 4.25, 5.05, 0.55, sz=10.5, color=DARK_GRAY)
tb(slide, "Reward improvement:  +163%", 7.45, BODY_TOP + 5.0, 5.3, 0.35,
   sz=11.5, bold=True, color=STEEL_BLUE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Contribution 2: PPO vs DPO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "Contribution 2 \u2014 PPO vs. DPO Comparative Analysis", 6)
BODY_TOP = HEADER_H + 0.14

section_label(slide, "C2 Overview", 0.5, BODY_TOP)
tb(slide,
   "Head-to-head comparison: identical dataset (IMDB), base model (GPT-2), hardware, and evaluation protocol.",
   0.5, BODY_TOP + 0.28, 12.4, 0.42, sz=12.5, italic=True, color=MID_GRAY)

thin_rule(slide, 0.5, BODY_TOP + 0.78, 12.4)

# Comparison table
table_data = [
    ["Dimension",           "PPO",                             "DPO"],
    ["Pipeline stages",     "3  (SFT  \u2192  RM  \u2192  PPO)",  "2  (SFT  \u2192  DPO)"],
    ["Reward model",        "Separate RM required",            "No RM needed"],
    ["Training mode",       "Online (generates new samples)",  "Offline (preference pairs)"],
    ["KL regularisation",   "Explicit beta coefficient",       "Closed-form via B\u2013T model"],
    ["Expected speed",      "Baseline (reference)",            "~10x faster (Rafailov 2023)"],
    ["Expected reward",     "Higher peak reward",              "Within 5% of PPO peak"],
    ["Collapse rate",       "~20% across seeds",               "~2% across seeds"],
]
add_table(slide, table_data,
          col_widths=[2.8, 4.55, 4.55],
          l=0.5, t=BODY_TOP + 0.95,
          header_fill=NAVY, body_text_sz=11.5, header_text_sz=12)

# Measurement callout
rect(slide, 0.5, BODY_TOP + 5.05, 12.4, 0.85, fill=WARM_WHITE, line=RULE_GRAY)
left_accent(slide, 0.5, BODY_TOP + 5.05, 0.85, color=GOLD)
section_label(slide, "Metrics measured for both methods", 0.65, BODY_TOP + 5.1)
tb(slide,
   "Mean reward score (mean \u00b1 std)  \u2022  Perplexity on held-out reviews  \u2022  Wall-clock training time  \u2022  Collapse rate (3 seeds)  \u2022  Qualitative response quality (10 prompts)",
   0.65, BODY_TOP + 5.35, 12.1, 0.45, sz=11.5, color=CHARCOAL)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Contribution 3: Stability & Alignment Tax
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "Contribution 3 \u2014 Stability Benchmarking & Alignment Tax", 7)
BODY_TOP = HEADER_H + 0.14

section_label(slide, "C3 Overview", 0.5, BODY_TOP)
tb(slide,
   "Alignment Tax: the increase in perplexity (degradation of general language modelling ability) "
   "caused by RLHF fine-tuning. Measured across KL coefficient values and training seeds.",
   0.5, BODY_TOP + 0.28, 12.4, 0.55, sz=12.5, color=CHARCOAL)

thin_rule(slide, 0.5, BODY_TOP + 0.92, 12.4)

# Beta comparison table
beta_data = [
    ["KL Coefficient (beta)",  "Training Stability",  "Reward Gain",   "Alignment Tax",  "Expected Collapse Rate"],
    ["0.05",                   "Low",                 "Very High",     "High",           "~40%"],
    ["0.10",                   "Moderate",            "High",          "Moderate",       "~15%"],
    ["0.20  (default)",        "High",                "Good",          "Low",            "< 5%"],
    ["0.50",                   "Very High",           "Limited",       "Very Low",       "< 1%"],
]
add_table(slide, beta_data,
          col_widths=[2.7, 2.4, 2.1, 2.3, 2.4],
          l=0.5, t=BODY_TOP + 1.05,
          header_fill=CHARCOAL, body_text_sz=11.5, header_text_sz=12)

thin_rule(slide, 0.5, BODY_TOP + 3.2, 12.4)

# Left: Measurements
section_label(slide, "Measurements  (5 seeds per beta)", 0.5, BODY_TOP + 3.3)
bullets(slide, [
    "Collapse rate: % of runs with reward hacking onset",
    "KL divergence trajectory: optimal stopping threshold",
    "Perplexity delta: quantified alignment tax",
    "Reward vs. perplexity trade-off curve",
    "Reward normalisation impact (Zheng et al., 2023)",
], 0.5, BODY_TOP + 3.65, 6.3, 2.6, sz=12.5)

# Right: Expected findings
rect(slide, 7.0, BODY_TOP + 3.3, 5.9, 2.9, fill=WARM_WHITE, line=RULE_GRAY)
left_accent(slide, 7.0, BODY_TOP + 3.3, 2.9, color=TEAL)
section_label(slide, "Expected Findings", 7.15, BODY_TOP + 3.35)
bullets(slide, [
    "Stable training window: beta in {0.1, 0.2}",
    "Reward hacking: onset at KL ~10-15 nats",
    "Alignment tax <= 15% perplexity increase",
    "DPO shows lower variance across seeds than PPO",
    "Reward normalisation reduces collapse by ~18%",
], 7.15, BODY_TOP + 3.65, 5.6, 2.4, sz=12, color=CHARCOAL)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Methodology
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "Methodology", 8)
BODY_TOP = HEADER_H + 0.1

# Dataset card
rect(slide, 0.45, BODY_TOP, 3.8, 2.9, fill=WARM_WHITE, line=RULE_GRAY)
left_accent(slide, 0.45, BODY_TOP, 2.9, color=STEEL_BLUE)
section_label(slide, "Dataset", 0.58, BODY_TOP + 0.06)
bullets(slide, [
    "IMDB Movie Reviews (50,000)",
    "Positive reviews  \u2192  SFT corpus",
    "Pos/neg pairs  \u2192  RM training",
    "Test prefixes  \u2192  PPO/DPO prompts",
], 0.58, BODY_TOP + 0.35, 3.55, 2.35, sz=12)

# Model stack card
rect(slide, 4.55, BODY_TOP, 3.8, 2.9, fill=WARM_WHITE, line=RULE_GRAY)
left_accent(slide, 4.55, BODY_TOP, 2.9, color=TEAL)
section_label(slide, "Model Stack", 4.68, BODY_TOP + 0.06)
bullets(slide, [
    "Base: GPT-2 (124M parameters)",
    "RM: GPT-2 + classification head",
    "Policy: GPT-2 (SFT init)",
    "Reference: frozen SFT (KL target)",
    "Library: HuggingFace TRL",
], 4.68, BODY_TOP + 0.35, 3.55, 2.35, sz=12)

# Config card
rect(slide, 8.65, BODY_TOP, 4.65, 2.9, fill=WARM_WHITE, line=RULE_GRAY)
left_accent(slide, 8.65, BODY_TOP, 2.9, color=GOLD)
section_label(slide, "Hyperparameters", 8.78, BODY_TOP + 0.06)
bullets(slide, [
    "PPO LR: 1.41e-5  |  DPO LR: 1e-5",
    "PPO batch: 16  |  DPO batch: 32",
    "PPO steps: 200  |  DPO epochs: 3",
    "Max gen. tokens: 32",
    "KL beta (PPO): 0.2  (swept: 0.05-0.5)",
    "Seeds per config: 3-5",
], 8.78, BODY_TOP + 0.35, 4.4, 2.35, sz=12)

thin_rule(slide, 0.45, BODY_TOP + 3.07, 12.85)

# Evaluation metrics
section_label(slide, "Evaluation Metrics", 0.5, BODY_TOP + 3.15)
metrics = [
    ("Mean Reward Score",        "Primary alignment quality measure"),
    ("Reward Std Deviation",     "Stability across seeds"),
    ("KL Divergence",            "Policy drift from reference"),
    ("Perplexity (PPL)",         "Alignment tax indicator"),
    ("Collapse Rate (%)",        "Reward hacking incidence"),
    ("Wall-clock Training Time", "Computational cost"),
]
for idx, (m_title, m_desc) in enumerate(metrics):
    col = idx % 3
    row = idx // 3
    lx = 0.45 + col * 4.3
    ty = BODY_TOP + 3.5 + row * 1.2
    rect(slide, lx, ty, 4.0, 1.05, fill=WARM_WHITE, line=RULE_GRAY)
    left_accent(slide, lx, ty, 1.05, color=STEEL_BLUE)
    tb(slide, m_title, lx + 0.13, ty + 0.08, 3.75, 0.42, sz=11.5, bold=True, color=NAVY)
    tb(slide, m_desc,  lx + 0.13, ty + 0.55, 3.75, 0.38, sz=10.5, color=CHARCOAL)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Implementation Timeline
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "10-Day Implementation Timeline", 9)
BODY_TOP = HEADER_H + 0.14

days_data = [
    ("Days 1\u20132", "C1: SFT + RM + PPO Pipeline", NAVY, [
        "Install TRL stack & prepare IMDB dataset",
        "Fine-tune SFT checkpoint",
        "Train Reward Model (GPT-2 classifier)",
        "Run 200 PPO training steps",
        "Log reward curve & KL trajectory",
    ]),
    ("Days 3\u20134", "C2: DPO Comparison", CHARCOAL, [
        "Implement DPO training loop",
        "Run on identical IMDB task",
        "Head-to-head evaluation protocol",
        "Compare reward & perplexity",
        "Wall-clock training time benchmark",
    ]),
    ("Days 5\u20136", "C3: Stability Benchmark", TEAL, [
        "5-seed repeatability runs (PPO & DPO)",
        "Sweep beta in {0.05, 0.1, 0.2, 0.5}",
        "Measure collapse rate per config",
        "KL divergence trajectory analysis",
        "Alignment tax perplexity curve",
    ]),
    ("Days 7\u20138", "Analysis & Reporting", STEEL_BLUE, [
        "Compile all metrics into tables",
        "Generate learning curve plots",
        "Draft results & discussion sections",
        "Qualitative response analysis",
        "Comparison figure (PPO vs DPO)",
    ]),
    ("Days 9\u201310", "Finalise & Submit", GOLD, [
        "Polish interactive CLI demo",
        "Write README and reproduce guide",
        "Proofread and format report",
        "Code review & clean-up",
        "Final submission",
    ]),
]
for i, (day, phase, color, tasks) in enumerate(days_data):
    lx = 0.3 + i * 2.6
    rect(slide, lx, BODY_TOP, 2.42, 0.82, fill=color)
    tb(slide, day,   lx + 0.08, BODY_TOP + 0.05, 2.3, 0.35, sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, phase, lx + 0.08, BODY_TOP + 0.42, 2.3, 0.38, sz=9.5, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, lx, BODY_TOP + 0.82, 2.42, 5.45, fill=WARM_WHITE, line=RULE_GRAY)
    rect(slide, lx, BODY_TOP + 0.82, 0.04, 5.45, fill=color)
    for j, task in enumerate(tasks):
        tb(slide, "\u2022  " + task, lx + 0.1, BODY_TOP + 1.0 + j * 1.0, 2.25, 0.88, sz=10.5, color=CHARCOAL)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Expected Results
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "Expected Results", 10)
BODY_TOP = HEADER_H + 0.14

section_label(slide, "Quantitative Targets", 0.5, BODY_TOP)

results_table = [
    ["Metric",                      "C1  PPO Pipeline",  "C2  DPO vs PPO",          "C3  Stability"],
    ["Reward improvement",          "+40\u201370% after 200 steps", "DPO within 5% of PPO",  "Measured per beta"],
    ["Training time",               "< 2 hrs (single GPU)", "DPO ~10x faster",        "5 seeds x 4 betas"],
    ["Perplexity increase (PPL)",   "< 15% (alignment tax)", "DPO lower PPL increase", "<= 15% at beta=0.2"],
    ["Collapse rate",               "< 5% with reward norm", "DPO ~2%  vs  PPO ~20%", "40% at beta=0.05"],
    ["KL divergence at stability",  "10\u201315 nats",       "Lower KL for DPO",        "Threshold identified"],
]
add_table(slide, results_table,
          col_widths=[3.0, 3.1, 3.1, 3.1],
          l=0.45, t=BODY_TOP + 0.35,
          header_fill=NAVY, body_text_sz=11.5, header_text_sz=12)

thin_rule(slide, 0.45, BODY_TOP + 3.65, 12.4)

section_label(slide, "Qualitative Output Example", 0.5, BODY_TOP + 3.75)
rect(slide, 0.45, BODY_TOP + 4.05, 5.85, 1.95, fill=RGBColor(0xF9, 0xEB, 0xEB), line=RULE_GRAY)
left_accent(slide, 0.45, BODY_TOP + 4.05, 1.95, color=RED)
tb(slide, "Before PPO  (SFT only)", 0.6, BODY_TOP + 4.1, 5.55, 0.42, sz=11.5, bold=True, color=RED)
tb(slide, "\"...slow and tedious. The plot dragged on without purpose. I would not recommend this film.\"",
   0.6, BODY_TOP + 4.55, 5.55, 0.8, sz=11, italic=True, color=CHARCOAL)
tb(slide, "Reward: \u22120.82", 0.6, BODY_TOP + 5.55, 5.55, 0.35, sz=11, bold=True, color=RED)

rect(slide, 6.95, BODY_TOP + 4.05, 5.85, 1.95, fill=RGBColor(0xEA, 0xF6, 0xEC), line=RULE_GRAY)
left_accent(slide, 6.95, BODY_TOP + 4.05, 1.95, color=GREEN)
tb(slide, "After PPO  (RLHF aligned)", 7.1, BODY_TOP + 4.1, 5.55, 0.42, sz=11.5, bold=True, color=GREEN)
tb(slide, "\"...a genuine masterpiece of storytelling. The performances are outstanding and the narrative grips you throughout.\"",
   7.1, BODY_TOP + 4.55, 5.55, 0.8, sz=11, italic=True, color=CHARCOAL)
tb(slide, "Reward: +2.14", 7.1, BODY_TOP + 5.55, 5.55, 0.35, sz=11, bold=True, color=GREEN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Future Work
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "Future Work Extensions", 11)
BODY_TOP = HEADER_H + 0.14

# C4 card
rect(slide, 0.45, BODY_TOP, 6.1, 5.6, fill=WARM_WHITE, line=RULE_GRAY)
left_accent(slide, 0.45, BODY_TOP, 5.6, color=TEAL)
# Badge
rect(slide, 0.55, BODY_TOP + 0.1, 0.72, 0.42, fill=TEAL)
tb(slide, "C4", 0.56, BODY_TOP + 0.11, 0.70, 0.40, sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(slide, "Human-Likeness Investigation",
   1.35, BODY_TOP + 0.1, 4.95, 0.42, sz=13.5, bold=True, color=NAVY)
thin_rule(slide, 0.55, BODY_TOP + 0.62, 5.85, color=RULE_GRAY)
bullets(slide, [
    "Inspired by Binz et al. (2025): cognitive model alignment",
    "Dataset: Psych-201 behavioural task battery",
    "Test if PPO-aligned GPT-2 better mimics human\ncognitive and decision-making patterns",
    "Metrics: behavioural alignment score, risk aversion\n  index, probability weighting function",
    "Hypothesis: RLHF narrows the gap to human\n  decision-making profiles",
], 0.58, BODY_TOP + 0.75, 5.85, 4.65, sz=12.5, color=CHARCOAL)

# C5 card
rect(slide, 6.85, BODY_TOP, 6.1, 5.6, fill=WARM_WHITE, line=RULE_GRAY)
left_accent(slide, 6.85, BODY_TOP, 5.6, color=GOLD)
# Badge
rect(slide, 6.95, BODY_TOP + 0.1, 0.72, 0.42, fill=GOLD)
tb(slide, "C5", 6.96, BODY_TOP + 0.11, 0.70, 0.40, sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(slide, "ADPO: Active Preference Learning",
   7.75, BODY_TOP + 0.1, 4.95, 0.42, sz=13.5, bold=True, color=NAVY)
thin_rule(slide, 6.95, BODY_TOP + 0.62, 5.85, color=RULE_GRAY)
bullets(slide, [
    "Active DPO: selectively query RM for uncertain pairs",
    "Based on: Liu et al. (2024) active preference framework",
    "Query strategy: uncertainty sampling via reward\n  model confidence scores",
    "Goal: reduce annotation cost by 40-60%\n  while preserving alignment quality",
    "Comparison: ADPO vs full DPO vs random sampling\n  at matched label budgets",
], 6.98, BODY_TOP + 0.75, 5.85, 4.65, sz=12.5, color=CHARCOAL)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — References (two-column, all 17)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_slide_chrome(slide, "References  (APA 7.0)", 12)
BODY_TOP = HEADER_H + 0.06

all_refs = [
    "[1]  Ziegler et al. (2019). Fine-tuning language models from human preferences. arXiv:1909.08593",
    "[2]  Stiennon et al. (2020). Learning to summarize with human feedback. NeurIPS 33, 3008-3021.",
    "[3]  Christiano et al. (2017). Deep RL from human preferences. NeurIPS 30.",
    "[4]  Schulman et al. (2017). Proximal policy optimization algorithms. arXiv:1707.06347",
    "[5]  Brown et al. (2020). Language models are few-shot learners. NeurIPS 33, 1877-1901.",
    "[6]  Ouyang et al. (2022). Training LMs to follow instructions (InstructGPT). NeurIPS 35.",
    "[7]  Bai et al. (2022). Training a helpful and harmless assistant with RLHF. arXiv:2204.05862",
    "[8]  Gao et al. (2023). Scaling laws for reward model overoptimization. ICML 2023.",
    "[9]  Rafailov et al. (2023). Direct preference optimization (DPO). NeurIPS 36.",
    "[10] Touvron et al. (2023). Llama 2: Open foundation and fine-tuned chat models. arXiv:2307.09288",
    "[11] Zheng et al. (2023). Secrets of RLHF in LLMs Part I: PPO. arXiv:2307.04964",
    "[12] Yuan et al. (2023). RRHF: Rank responses to align LMs with human feedback. NeurIPS 36.",
    "[13] Dong et al. (2023). RAFT: Reward ranked finetuning for generative foundation model alignment. TMLR.",
    "[14] Azar et al. (2024). A general theoretical paradigm to understand learning from human feedback. AISTATS.",
    "[15] Wang et al. (2024). A comprehensive survey of LLM alignment: RLHF, RLAIF, PPO, DPO. arXiv:2407.16216",
    "[16] Binz et al. (2025). Turning large language models into cognitive models. arXiv:2306.03917",
    "[17] Liu et al. (2024). Active direct preference optimization. arXiv:2402.10141",
]
mid = (len(all_refs) + 1) // 2  # 9 left, 8 right
left_refs  = all_refs[:mid]
right_refs = all_refs[mid:]

for i, ref in enumerate(left_refs):
    ty = BODY_TOP + 0.12 + i * 0.67
    rect(slide, 0.45, ty, 0.04, 0.52, fill=STEEL_BLUE)
    tb(slide, ref, 0.6, ty, 6.0, 0.62, sz=9.5, color=CHARCOAL)

for i, ref in enumerate(right_refs):
    ty = BODY_TOP + 0.12 + i * 0.67
    rect(slide, 6.85, ty, 0.04, 0.52, fill=STEEL_BLUE)
    tb(slide, ref, 7.0, ty, 6.0, 0.62, sz=9.5, color=CHARCOAL)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Closing / Q&A
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
rect(slide, 0, 0, SLIDE_W, 0.06, fill=GOLD)
rect(slide, 0, SLIDE_H - 0.06, SLIDE_W, 0.06, fill=GOLD)

# Decorative horizontal rule
rect(slide, 1.5, 3.55, 10.33, 0.05, fill=GOLD)

tb(slide, "Thank You", 1, 1.4, 11.33, 1.4,
   sz=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(slide, "Questions & Discussion",
   1, 2.9, 11.33, 0.7, sz=20, italic=True, color=SILVER, align=PP_ALIGN.CENTER)

tb(slide,
   "C1  Reproducible RLHF Pipeline     \u2022     C2  PPO vs DPO Comparison     \u2022     C3  Stability & Alignment Tax Benchmarking",
   1, 3.75, 11.33, 0.55, sz=12, color=MID_GRAY, align=PP_ALIGN.CENTER)
tb(slide,
   "Future Work:   C4  Human-Likeness Investigation   \u2022   C5  ADPO Active Preference Learning",
   1, 4.35, 11.33, 0.5, sz=11.5, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

thin_rule(slide, 3.5, 5.1, 6.33, color=CHARCOAL)
tb(slide, "[Student Name]  |  Masters in AI  |  [University]  |  2025",
   1, 5.2, 11.33, 0.5, sz=12, color=MID_GRAY, align=PP_ALIGN.CENTER)


ppt_path = "RLHF_Project_Proposal.pptx"
prs.save(ppt_path)
print(f"✅  Saved: {ppt_path}")
print("\nDone! Both files created in the current directory.")
